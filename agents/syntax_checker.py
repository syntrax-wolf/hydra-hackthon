"""Syntax Checker Agent — validates generated code before execution.

Two-phase checking:
1. Python AST parsing for basic syntax errors (fast, local)
2. LLM-based deep API review for library-specific mistakes (thorough)
"""

import ast
import json
import re
import logging
import time
import httpx
from core.config import config

log = logging.getLogger("syntax_checker")

MAX_RETRIES = 2  # Max retry attempts for LLM calls (total attempts = MAX_RETRIES + 1)
MAX_CODE_LENGTH = 15000

SYNTAX_CHECKER_SYSTEM_PROMPT = """You are an expert Python syntax reviewer, API specialist, and visual design quality auditor. You check code for correctness AND visual quality. Your job has TWO parts:

PART A — Find syntax errors, API misuse, and incorrect method calls.
PART B — Check that visual design standards are followed (colors, spacing, styling).

═══════════════════════════════════════════════
PART A: API CORRECTNESS RULES
═══════════════════════════════════════════════

PYTHON-PPTX (pptx):
  CORRECT IMPORTS:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu, Cm
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE

  MISTAKES TO CATCH:
    ✗ RGBColor('#1a73e8')          → ✓ RGBColor(0x1a, 0x73, 0xe8) [3 ints, NOT hex string]
    ✗ RGBColor('1a73e8')           → ✓ RGBColor(0x1a, 0x73, 0xe8)
    ✗ paragraph.font.size = 14     → ✓ paragraph.font.size = Pt(14)
    ✗ paragraph.font.color = ...   → ✓ paragraph.font.color.rgb = RGBColor(...)
    ✗ shape.left = 1               → ✓ shape.left = Inches(1)
    ✗ prs.slide_layouts['Title']   → ✓ prs.slide_layouts[0] [index only]
    ✗ prs.add_slide(layout)        → ✓ prs.slides.add_slide(layout)
    ✗ from pptx.util import Points → ✓ from pptx.util import Pt
    ✗ paragraph.alignment = 'center' → ✓ paragraph.alignment = PP_ALIGN.CENTER
    ✗ cell.text = number           → ✓ cell.text = str(number) [must be string]
    ✗ add_table(rows, cols)        → needs left, top, width, height args too
    ✗ slide.placeholders[...]      → NEVER use placeholders, use add_textbox()

REPORTLAB:
  CORRECT IMPORTS:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.units import inch, cm, mm
    from reportlab.lib.colors import HexColor, Color, black, white
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, PageBreak

  MISTAKES TO CATCH:
    ✗ HexColor('1a73e8')           → ✓ HexColor('#1a73e8') [NEEDS # prefix]
    ✗ from reportlab... import RGBColor → ✓ Use HexColor or Color(r,g,b)
    ✗ SimpleDocTemplate('f.pdf')   → ✓ needs pagesize=letter
    ✗ Paragraph(text)              → ✓ Paragraph(text, style) [style required]
    ✗ Spacer(12)                   → ✓ Spacer(1, 12) [needs width AND height]
    ✗ Table(data, colWidths=100)   → ✓ colWidths must be a list
    ✗ Image('path', 400, 300)      → ✓ Image('path', width=400, height=300)
    ✗ elements.append("string")    → ✓ must be a flowable (Paragraph, etc.)

PYTHON-DOCX (docx):
    ✗ Document('new.docx')         → ✓ Document() [no arg for new doc]
    ✗ run.font.color = RGBColor    → ✓ run.font.color.rgb = RGBColor(...)
    ✗ RGBColor('#hex')             → ✓ RGBColor(0xRR, 0xGG, 0xBB) [3 ints]
    ✗ from docx.shared import Points → ✓ from docx.shared import Pt

MATPLOTLIB:
    ✗ Missing matplotlib.use('Agg') before pyplot import/figure creation
    ✗ plt.show() in headless mode  → remove it
    ✗ Missing plt.close() after savefig → memory leak
    ✗ Variable named `colors`      → shadows reportlab import, rename it

PIL/PILLOW:
    ✗ ImageFont.truetype() without try/except fallback → MUST have fallback to load_default()
    ✗ PILImage.new('RGBA', ...) saved to JPEG → RGBA not supported in JPEG, convert to RGB
    ✗ draw.rounded_rectangle() without Pillow >= 8.2 check → use draw.rectangle as fallback
    ✗ Font size too small (< 12px) → text will be unreadable

GENERAL PYTHON:
    ✗ Unterminated f-strings: f"...{e) → f"...{e}"
    ✗ Missing imports
    ✗ String/number type mismatches
    ✗ Variables used before assignment
    ✗ Unclosed parentheses/brackets/strings
    ✗ Division by zero (check if denominator can be 0)

═══════════════════════════════════════════════
PART B: VISUAL DESIGN QUALITY CHECKS
═══════════════════════════════════════════════

Beyond syntax, check these DESIGN issues (report as severity: "warning"):

1. COLOR QUALITY:
   ✗ Using pure black (#000000) for text → should use #1a202c or #2d3748
   ✗ Using pure white (#ffffff) for slide backgrounds → should use #f7fafc or #fafbfc
   ✗ Using random/ugly colors for charts → should use a coordinated palette
   ✗ Red for positive values or green for negative → semantically wrong
   ✗ No color variation (everything is one color) → needs a palette
   ✗ Low contrast text on background (e.g., light gray on white)

2. CHART QUALITY:
   ✗ No charts in a document with data → MUST have 2-3 charts
   ✗ Charts missing titles or axis labels
   ✗ Charts with plt.show() instead of savefig
   ✗ Charts saved without plt.close() → memory leak
   ✗ No value labels on bar charts → add ax.bar_label()
   ✗ Chart background doesn't match document theme
   ✗ Missing ax.spines removal (top, right spines should be hidden)
   ✗ No grid lines on y-axis for bar/line charts

3. TABLE QUALITY:
   ✗ Tables without header styling (no background color, no bold)
   ✗ Tables without alternating row colors
   ✗ Tables without cell padding
   ✗ Numbers not right-aligned in table cells
   ✗ No border/grid on table

4. TYPOGRAPHY:
   ✗ Title text too small (PPTX title < Pt(28), PDF title < 20pt)
   ✗ Body text too small (PPTX < Pt(12), PDF < 9pt)
   ✗ No font size hierarchy (everything same size)
   ✗ Missing bold on headers/titles

5. SPACING & LAYOUT:
   ✗ No spacing between elements (no Spacer in PDF, no gap in PPTX)
   ✗ Content too close to slide/page edges (< 0.5 inches margin)
   ✗ Images without explicit dimensions (will render wrong size)

6. MISSING VISUAL ELEMENTS:
   ✗ No metric cards / KPI visuals when metric data is available
   ✗ No title page / cover page
   ✗ No section dividers between major sections
   ✗ No page numbers in PDF footer

═══════════════════════════════════════════════
RESPONSE FORMAT
═══════════════════════════════════════════════
{
  "has_errors": true/false,
  "errors": [
    {
      "line": "approximate line number or code snippet",
      "error_type": "syntax|api_misuse|import|type_mismatch|missing_import|logic|design_quality",
      "description": "What is wrong",
      "fix": "The corrected code for that section",
      "severity": "critical|warning"
    }
  ],
  "fixed_code": "the complete corrected Python code if errors were found, or empty string if no errors",
  "summary": "Brief summary of findings"
}

RULES:
1. ONLY report REAL errors. Do not flag correct code.
2. If code is correct, return has_errors: false with empty errors.
3. For critical errors, provide fixed_code with ALL corrections.
4. Preserve original logic — fix errors, don't refactor.
5. RGBColor usage is the #1 bug — check EVERY instance.
6. Design issues are severity "warning" (not "critical") unless they cause runtime errors.
7. When providing fixed_code, also apply design improvements to make the document beautiful.
8. If the code generates a document with NO charts despite having data, mark as critical (charts are mandatory)."""


async def _call_llm(system_prompt: str, user_message: str, label: str = "SYNTAX") -> str:
    """Call OpenRouter for syntax checking."""
    last_error = None
    for attempt in range(MAX_RETRIES + 1):
        try:
            log.info("[%s] Calling OpenRouter (attempt=%d/%d)...", label, attempt + 1, MAX_RETRIES + 1)
            t0 = time.time()
            async with httpx.AsyncClient(timeout=120) as client:
                response = await client.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {config.openrouter_api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": config.openrouter_model,
                        "messages": [
                            {"role": "system", "content": system_prompt},
                            {"role": "user", "content": user_message},
                        ],
                        "max_tokens": 8192,
                        "temperature": 0.1,
                        "top_p": 0.95,
                        "reasoning": {"effort": "high"},
                    },
                )
                elapsed = time.time() - t0
                response.raise_for_status()
                data = response.json()
                content = data["choices"][0]["message"]["content"]

                # Strip <think> blocks
                content = re.sub(r"<think>.*?</think>", "", content, flags=re.DOTALL).strip()
                # Strip markdown fences
                if content.startswith("```"):
                    content = content.split("\n", 1)[1] if "\n" in content else content[3:]
                if content.endswith("```"):
                    content = content[:-3]
                # Strip pre-JSON text
                content = content.strip()
                if content and not content.startswith("{"):
                    idx = content.find("{")
                    if idx > 0:
                        content = content[idx:]

                log.info("[%s] Response in %.1fs", label, elapsed)
                return content.strip()
        except (httpx.HTTPStatusError, httpx.ConnectError, httpx.ReadTimeout) as e:
            last_error = e
            log.warning("[%s] Attempt %d failed: %s", label, attempt + 1, e)
            if attempt < MAX_RETRIES:
                import asyncio
                await asyncio.sleep(2)
    raise last_error


def _parse_json(raw: str) -> dict:
    """Parse JSON from LLM response."""
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        pass
    match = re.search(r"\{.*\}", raw, re.DOTALL)
    if match:
        try:
            return json.loads(match.group())
        except json.JSONDecodeError:
            pass
    return {"has_errors": False, "errors": [], "fixed_code": "", "summary": "LLM response unparseable — proceeding with original code"}


def check_ast(code: str) -> dict:
    """Phase 1: Quick Python AST parse to catch basic syntax errors."""
    try:
        ast.parse(code)
        return {"has_errors": False, "errors": [], "fixed_code": "", "summary": "AST parse OK"}
    except SyntaxError as e:
        return {
            "has_errors": True,
            "errors": [{
                "line": str(e.lineno or "unknown"),
                "error_type": "syntax",
                "description": f"SyntaxError: {e.msg} (line {e.lineno}, col {e.offset})",
                "fix": "",
                "severity": "critical",
            }],
            "fixed_code": "",
            "summary": f"Python syntax error at line {e.lineno}: {e.msg}",
        }


async def check(code: str) -> dict:
    """Check code for syntax errors and API misuse.

    Phase 1: AST parse (fast, local).
    Phase 2: LLM deep API review (catches library-specific mistakes).

    Returns dict with: has_errors, errors, fixed_code, summary
    """
    if not code or not code.strip():
        return {
            "has_errors": True,
            "errors": [{"line": "0", "error_type": "syntax", "description": "Empty code", "fix": "", "severity": "critical"}],
            "fixed_code": "",
            "summary": "Code is empty",
        }

    # Phase 1: Quick AST syntax check
    ast_result = check_ast(code)
    if ast_result["has_errors"]:
        log.info("[SYNTAX] AST error found: %s", ast_result["errors"][0]["description"])
        # Ask LLM to fix syntax errors
        error_desc = "\n".join(f"- {e['description']}" for e in ast_result["errors"])
        truncated = code[:MAX_CODE_LENGTH]
        raw = await _call_llm(
            SYNTAX_CHECKER_SYSTEM_PROMPT,
            f"This Python code has syntax errors detected by the AST parser:\n\nERRORS:\n{error_desc}\n\nCODE:\n```python\n{truncated}\n```\n\nFix ALL syntax errors and also check for any API misuse in document-generation libraries.",
            label="SYNTAX-FIX",
        )
        result = _parse_json(raw)
        result.setdefault("has_errors", True)
        result.setdefault("errors", ast_result["errors"])
        result.setdefault("fixed_code", "")
        result.setdefault("summary", "Syntax errors found")
        return result

    # Phase 2: Deep LLM-based API review
    truncated = code[:MAX_CODE_LENGTH] if len(code) > MAX_CODE_LENGTH else code
    log.info("[SYNTAX] AST OK — running deep LLM API review (%d chars)", len(truncated))
    raw = await _call_llm(
        SYNTAX_CHECKER_SYSTEM_PROMPT,
        f"Review this Python code for syntax errors and API misuse. Focus especially on python-pptx, reportlab, python-docx, and matplotlib APIs.\n\nCODE:\n```python\n{truncated}\n```",
        label="SYNTAX-REVIEW",
    )
    result = _parse_json(raw)
    result.setdefault("has_errors", False)
    result.setdefault("errors", [])
    result.setdefault("fixed_code", "")
    result.setdefault("summary", "No issues found")
    return result
